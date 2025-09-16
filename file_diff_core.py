# -*- coding: utf-8 -*-
#!/usr/bin/env python
"""
文件内容对比 - 核心功能模块

这个模块包含所有与UI无关的核心功能实现，包括：
- 各种文件格式解析（PDF、Word、XLSX、PPT、Markdown等）
- 文件内容提取
- 文件哈希计算
- 文件对比算法
- 差异分析与报告生成

这些功能可以被UI模块调用，也可以作为独立库使用。
"""

from __future__ import annotations
import os
import hashlib
import shutil
import difflib
import re
from typing import List, Dict, Tuple, Optional, Any, Set, Iterable
import mimetypes

# 尝试导入各种文件格式处理库
try:
    import PyPDF2
    _has_pypdf2 = True
except ImportError:
    _has_pypdf2 = False

try:
    import docx
    _has_docx = True
except ImportError:
    _has_docx = False
try:
    import pandas as pd
    _has_pandas = True
except ImportError:
    _has_pandas = False

try:
    import markdown
    _has_markdown = True
except ImportError:
    _has_markdown = False

try:
    from pptx import Presentation
    _has_pptx = True
except ImportError:
    _has_pptx = False

# 支持的文件格式
supported_formats = {
    '.pdf': 'PDF文档',
    '.docx': 'Word文档',
    '.doc': 'Word文档(旧版)',
    '.xlsx': 'Excel表格',
    '.xls': 'Excel表格(旧版)',
    '.pptx': 'PowerPoint演示文稿',
    '.md': 'Markdown文档',
    '.txt': '文本文件',
    '.py': 'Python源代码',
    '.html': 'HTML文档',
    '.css': 'CSS样式表',
    '.js': 'JavaScript文件',
    '.bat': '批处理文件',
    '.cmd': '命令文件',
    '.sh': 'Shell脚本',
    '.ps1': 'PowerShell脚本',
    '.php': 'PHP脚本',
    '.java': 'Java源代码',
    '.cpp': 'C++源代码',
    '.c': 'C源代码',
    '.h': '头文件',
    '.cs': 'C#源代码',
    '.vb': 'Visual Basic代码',
    '.sql': 'SQL脚本',
    '.xml': 'XML文件',
    '.json': 'JSON文件',
}

class FileNotSupportedError(Exception):
    """不支持的文件格式异常"""
    pass

class FileReadError(Exception):
    """文件读取错误异常"""
    pass

def get_file_type(file_path: str) -> str:
    """
    获取文件类型
    
    Args:
        file_path: 文件路径
        
    Returns:
        str: 文件类型描述
    """
    _, ext = os.path.splitext(file_path.lower())
    if ext in supported_formats:
        return supported_formats[ext]
    
    # 尝试使用mimetypes模块判断
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        return mime_type
    
    return "未知文件类型"

def calculate_file_hash(file_path: str, hash_algorithm: str = 'md5') -> str:
    """
    计算文件的哈希值
    
    Args:
        file_path: 文件路径
        hash_algorithm: 哈希算法，支持'md5', 'sha1', 'sha256'
        
    Returns:
        str: 文件哈希值的十六进制表示
    """
    hash_func = getattr(hashlib, hash_algorithm)
    hash_obj = hash_func()
    
    # 分块读取文件以处理大文件
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(4096), b''):
            hash_obj.update(chunk)
    
    return hash_obj.hexdigest()

def extract_text_from_file(file_path: str) -> str:
    """
    从各种文件格式中提取文本内容
    
    Args:
        file_path: 文件路径
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        FileNotSupportedError: 不支持的文件格式
        FileReadError: 文件读取错误
    """
    _, ext = os.path.splitext(file_path.lower())
    
    try:
        if ext == '.pdf' and _has_pypdf2:
            return _extract_text_from_pdf(file_path)
        elif ext == '.docx' and _has_docx:
            return _extract_text_from_docx(file_path)
        elif ext in ['.xlsx', '.xls'] and _has_pandas:
            return _extract_text_from_excel(file_path)
        elif ext == '.pptx' and _has_pptx:
            return _extract_text_from_pptx(file_path)
        elif ext == '.md' and _has_markdown:
            return _extract_text_from_markdown(file_path)
        elif ext in ['.txt', '.py', '.html', '.css', '.js', '.bat', '.cmd', '.sh', '.ps1', '.php', '.java', '.cpp', '.c', '.h', '.cs', '.vb', '.sql', '.xml', '.json']:
            return _extract_text_from_text_file(file_path)
        else:
            raise FileNotSupportedError(f"不支持的文件格式: {ext}")
    except Exception as e:
        if isinstance(e, FileNotSupportedError):
            raise
        raise FileReadError(f"读取文件失败: {str(e)}")

def _extract_text_from_pdf(file_path: str) -> str:
    """从PDF文件中提取文本"""
    text = []
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text.append(page.extract_text())
    return "\n\n".join(text)

def _extract_text_from_docx(file_path: str) -> str:
    """从Word文档中提取文本"""
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text)
    return "\n\n".join(text)

def _extract_text_from_excel(file_path: str) -> str:
    """从Excel文件中提取文本"""
    try:
        df = pd.read_excel(file_path)
        return str(df)
    except Exception:
        # 如果是多工作表文件，尝试读取所有工作表
        text = []
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text.append(f"===== {sheet_name} =====")
            text.append(str(df))
        return "\n\n".join(text)

def _extract_text_from_pptx(file_path: str) -> str:
    """从PowerPoint文件中提取文本"""
    prs = Presentation(file_path)
    text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"幻灯片 {slide_num}:"]
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_text.append(shape.text)
        if len(slide_text) > 1:  # 如果幻灯片有内容
            text.append("\n".join(slide_text))
    return "\n\n".join(text)

def _extract_text_from_markdown(file_path: str) -> str:
    """从Markdown文件中提取文本"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    return content

def _extract_text_from_text_file(file_path: str) -> str:
    """从文本文件中提取文本"""
    # 尝试多种编码
    encodings = ['utf-8', 'gbk', 'latin-1']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    
    # 如果都失败了，返回错误信息
    raise FileReadError(f"无法使用任何支持的编码打开文件: {file_path}")

def compare_files(file1_path: str, file2_path: str, ignore_whitespace: bool = False) -> Dict[str, Any]:
    """
    对比两个文件的内容
    
    Args:
        file1_path: 第一个文件路径
        file2_path: 第二个文件路径
        ignore_whitespace: 是否忽略空白字符差异
        
    Returns:
        Dict: 对比结果，包含差异信息
    """
    try:
        # 提取两个文件的文本内容
        text1 = extract_text_from_file(file1_path)
        text2 = extract_text_from_file(file2_path)
        
        # 如果忽略空白字符，预处理文本
        if ignore_whitespace:
            text1 = re.sub(r'\s+', ' ', text1).strip()
            text2 = re.sub(r'\s+', ' ', text2).strip()
        
        # 计算文件哈希值
        hash1 = calculate_file_hash(file1_path)
        hash2 = calculate_file_hash(file2_path)
        
        # 判断文件是否完全相同
        if hash1 == hash2:
            return {
                'is_same': True,
                'hash1': hash1,
                'hash2': hash2,
                'differences': [],
                'text1': text1,
                'text2': text2,
                'error': None
            }
        
        # 使用difflib计算差异
        differ = difflib.Differ()
        diff_lines = list(differ.compare(text1.splitlines(keepends=True), text2.splitlines(keepends=True)))
        
        # 生成差异报告
        differences = []
        for i, line in enumerate(diff_lines):
            if line.startswith('- ') or line.startswith('+ ') or line.startswith('? '):
                differences.append({
                    'line_number': i,
                    'line_content': line,
                    'type': line[0]
                })
        
        # 计算相似度
        similarity = difflib.SequenceMatcher(None, text1, text2).ratio()
        
        return {
            'is_same': False,
            'hash1': hash1,
            'hash2': hash2,
            'differences': differences,
            'similarity': similarity,
            'text1': text1,
            'text2': text2,
            'error': None
        }
        
    except Exception as e:
        return {
            'is_same': False,
            'hash1': None,
            'hash2': None,
            'differences': [],
            'text1': None,
            'text2': None,
            'error': str(e)
        }

def iter_files(root: str, recursive: bool = True, skip_extensions: set = None) -> Iterable[str]:
    """
    遍历指定目录中的所有文件
    
    Args:
        root: 要扫描的根目录路径
        recursive: 是否递归扫描子目录
        skip_extensions: 要跳过的文件扩展名集合，如 {'.tmp', '.log'}
        
    Yields:
        str: 符合条件的文件完整路径
    """
    if skip_extensions is None:
        skip_extensions = set()
    
    for dirpath, dirs, files in os.walk(root):
        for f in files:
            filepath = os.path.join(dirpath, f)
            _, ext = os.path.splitext(f.lower())
            if ext not in skip_extensions:
                yield filepath
        
        # 如果不递归，只遍历顶层目录
        if not recursive:
            break

def compare_directories(dir1: str, dir2: str, recursive: bool = True) -> Dict[str, Any]:
    """
    对比两个目录的结构和内容
    
    Args:
        dir1: 第一个目录路径
        dir2: 第二个目录路径
        recursive: 是否递归对比子目录
        
    Returns:
        Dict: 目录对比结果
    """
    # 获取两个目录中的所有文件
    dir1_files = set()
    for file_path in iter_files(dir1, recursive):
        rel_path = os.path.relpath(file_path, dir1)
        dir1_files.add(rel_path)
    
    dir2_files = set()
    for file_path in iter_files(dir2, recursive):
        rel_path = os.path.relpath(file_path, dir2)
        dir2_files.add(rel_path)
    
    # 找出仅在第一个目录中存在的文件
    only_in_dir1 = dir1_files - dir2_files
    
    # 找出仅在第二个目录中存在的文件
    only_in_dir2 = dir2_files - dir1_files
    
    # 找出在两个目录中都存在的文件
    common_files = dir1_files & dir2_files
    
    # 对比共同文件的内容
    same_files = []
    different_files = []
    
    for rel_path in common_files:
        file1_path = os.path.join(dir1, rel_path)
        file2_path = os.path.join(dir2, rel_path)
        
        # 先比较文件大小
        size1 = os.path.getsize(file1_path)
        size2 = os.path.getsize(file2_path)
        
        if size1 != size2:
            different_files.append(rel_path)
            continue
        
        # 文件大小相同，再比较哈希值
        hash1 = calculate_file_hash(file1_path)
        hash2 = calculate_file_hash(file2_path)
        
        if hash1 == hash2:
            same_files.append(rel_path)
        else:
            different_files.append(rel_path)
    
    return {
        'only_in_dir1': list(only_in_dir1),
        'only_in_dir2': list(only_in_dir2),
        'same_files': same_files,
        'different_files': different_files,
        'total_files_dir1': len(dir1_files),
        'total_files_dir2': len(dir2_files),
    }

def generate_diff_report(compare_result: Dict[str, Any], file1_name: str, file2_name: str) -> str:
    """
    生成差异报告
    
    Args:
        compare_result: 文件对比结果
        file1_name: 第一个文件名
        file2_name: 第二个文件名
        
    Returns:
        str: 差异报告文本
    """
    report = []
    report.append("文件对比报告")
    report.append("=" * 50)
    report.append(f"文件1: {file1_name}")
    report.append(f"文件2: {file2_name}")
    report.append("\n")
    
    if compare_result['error']:
        report.append(f"错误: {compare_result['error']}")
        return "\n".join(report)
    
    if compare_result['is_same']:
        report.append("结论: 两个文件内容完全相同")
        report.append(f"哈希值: {compare_result['hash1']}")
    else:
        report.append("结论: 两个文件内容不同")
        report.append(f"文件1哈希值: {compare_result['hash1']}")
        report.append(f"文件2哈希值: {compare_result['hash2']}")
        
        if 'similarity' in compare_result:
            report.append(f"相似度: {compare_result['similarity']:.2%}")
        
        report.append("\n差异详情:")
        report.append("=" * 50)
        
        # 生成上下文差异
        if compare_result['text1'] and compare_result['text2']:
            diff = difflib.unified_diff(
                compare_result['text1'].splitlines(keepends=True),
                compare_result['text2'].splitlines(keepends=True),
                fromfile=file1_name,
                tofile=file2_name
            )
            report.append("\n".join(diff))
    
    return "\n".join(report)